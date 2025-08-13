import random
import requests
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from duckduckgo_search import DDGS
import google.generativeai as genai

# Initialize the client with your API key
genai.configure(api_key="Your api key")
model = genai.GenerativeModel("gemini-1.5-flash")
import json

BACKGROUND_COLORS = [
    RGBColor(255, 230, 230),  # Light Pink
    RGBColor(230, 255, 230),  # Light Green
    RGBColor(230, 230, 255),  # Light Blue
    RGBColor(255, 245, 230),  # Light Peach
    RGBColor(240, 255, 255),  # Light Cyan
    RGBColor(255, 255, 204),  # Light Yellow
]


def clean_model_output(text: str) -> str:
    # Remove code fences like ```json ... ```
    if text.startswith("```"):
        lines = text.splitlines()
        if lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].startswith("```"):
            lines = lines[:-1]
        text = "\n".join(lines)
    return text.strip()

def apply_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color
def generate_slide_content(topic):
    prompt = (
    f"Create a 6-slide presentation outline on the topic '{topic}'. "
    "Slides should be: 1) Title, 2) Overview, 3-5) Key points, 6) Conclusion. "
    "For each slide except title, provide a concise slide title and 5 bullet points. "
    "Return the output as JSON with keys: overview, key_points (list of 3 slides), and conclusion. "
    "Example format:\n"
    "{\n"
    "  \"overview\": [\"point1\", \"point2\", \"point3\", \"point4\", \"point5\"],\n"
    "  \"key_points\": [\n"
    "    {\"title\": \"Key Point 1\", \"points\": [\"p1\", \"p2\", \"p3\", \"p4\", \"p5\"]},\n"
    "    {\"title\": \"Key Point 2\", \"points\": [\"p1\", \"p2\", \"p3\", \"p4\", \"p5\"]},\n"
    "    {\"title\": \"Key Point 3\", \"points\": [\"p1\", \"p2\", \"p3\", \"p4\", \"p5\"]}\n"
    "  ],\n"
    "  \"conclusion\": [\"point1\", \"point2\", \"point3\", \"point4\", \"point5\"]\n"
    "}\n"
    "Make bullet points clear and concise."
    )

    response = model.generate_content(
        prompt,
        generation_config=genai.GenerationConfig(
            temperature=0.7,
            max_output_tokens=1000,
        )
    )
    text = clean_model_output(response.text)
    try:
        slides_data = json.loads(text)
        return slides_data
    except Exception as e:
        print("JSON parse error:", e)
        print("Raw model output:\n", text)
        return None

def create_presentation(topic, slides_data):
    prs = Presentation()

    # Slide 1: Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_background(slide, random.choice(BACKGROUND_COLORS))
    slide.shapes.title.text = topic
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"An overview of {topic}"

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide, random.choice(BACKGROUND_COLORS))
    slide.shapes.title.text = "Overview"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.word_wrap = True
    for point in slides_data.get('overview', []):
        p = tf.add_paragraph()
        p.text = point
        run = p.runs[0]
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = False
        run.font.name = 'Calibri'
        run.space_after = Pt(8)

    # Slides 3-5: Key points
    key_points = slides_data.get('key_points', [])
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        apply_background(slide, random.choice(BACKGROUND_COLORS))
        title = key_points[i]['title'] if i < len(key_points) else f"Key Point {i+1}"
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        points = key_points[i]['points'] if i < len(key_points) else []
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(6)

    # Slide 6: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide, random.choice(BACKGROUND_COLORS))
    slide.shapes.title.text = "Conclusion / Takeaways"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.word_wrap = True
    for point in slides_data.get('conclusion', []):
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(8)

    output_file = f"{topic.replace(' ', '_')}_presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation saved as {output_file}")

if __name__ == "__main__":
    topic = input("Enter your presentation topic: ").strip()
    if topic:
        slides_data = generate_slide_content(topic)
        if slides_data:
            create_presentation(topic, slides_data)
        else:
            print("❌ Failed to generate slide content.")
    else:
        print("⚠️ Topic cannot be empty.")
